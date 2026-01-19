# -*- coding: utf-8 -*-
"""The PowerPoint reader to read and chunk PowerPoint presentations."""
import base64
import hashlib
import json
from typing import Any, Literal

from ._reader_base import ReaderBase
from ._text_reader import TextReader
from ._utils import _get_media_type_from_data
from .._document import Document, DocMetadata
from ...message import ImageBlock, Base64Source, TextBlock
from ..._logging import logger


def _extract_table_data(table: Any) -> list[list[str]]:
    """Extract table data from a PowerPoint table.

    Args:
        table (`Any`):
            The table object from python-pptx.

    Returns:
        `list[list[str]]`:
            Table data represented as a 2D list, where each inner list
            represents a row, and each string in the row represents a cell.
    """
    table_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            # Extract text from cell, preserving line breaks within cells
            cell_text = cell.text.strip()
            # Replace line breaks with \n to preserve structure
            cell_text = cell_text.replace("\r\n", "\n").replace("\r", "\n")
            row_data.append(cell_text)
        table_data.append(row_data)
    return table_data


def _extract_images_from_shape(shape: Any) -> list[ImageBlock]:
    """Extract images from a shape (if it contains images).

    Args:
        shape (`Any`):
            The shape object from python-pptx.

    Returns:
        `list[ImageBlock]`:
            A list of ImageBlock objects, empty if no images found.
    """
    images = []

    # Check if shape is a picture
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        picture_type = MSO_SHAPE_TYPE.PICTURE
    except ImportError:
        picture_type = 13  # MSO_SHAPE_TYPE.PICTURE fallback

    if shape.shape_type == picture_type:
        try:
            # Get image data
            image_data = shape.image.blob

            # Determine media type
            media_type = _get_media_type_from_data(image_data)

            # Convert to base64
            base64_data = base64.b64encode(image_data).decode("utf-8")

            images.append(
                ImageBlock(
                    type="image",
                    source=Base64Source(
                        type="base64",
                        media_type=media_type,
                        data=base64_data,
                    ),
                ),
            )
        except Exception as e:
            logger.warning("Failed to extract image from shape: %s", e)

    return images


class PowerPointReader(ReaderBase):
    """The PowerPoint reader that supports reading text, image, and table
    content from PowerPoint presentations (.pptx files), and chunking the text
    content into smaller pieces.

    .. note:: The table content can be extracted in Markdown or JSON format.
    """

    def __init__(
        self,
        chunk_size: int = 512,
        split_by: Literal["char", "sentence", "paragraph"] = "sentence",
        include_image: bool = False,
        separate_slide: bool = False,
        separate_table: bool = False,
        table_format: Literal["markdown", "json"] = "markdown",
    ) -> None:
        """Initialize the PowerPoint reader.

        Args:
            chunk_size (`int`, default to 512):
                The size of each chunk, in number of characters.
            split_by (`Literal["char", "sentence", "paragraph"]`, default to \
            "sentence"):
                The unit to split the text, can be "char", "sentence", or
                "paragraph". The "sentence" option is implemented using the
                "nltk" library, which only supports English text.
            include_image (`bool`, default to False):
                Whether to include image content in the document. If True,
                images will be extracted and included as base64-encoded images.
            separate_slide (`bool`, default to False):
                Whether to treat each slide as a separate document. If True,
                each slide will be extracted as a separate Document object
                instead of being merged together.
            separate_table (`bool`, default to False):
                If True, tables will be treated as a new chunk to avoid
                truncation. But note when the table exceeds the chunk size,
                it will still be truncated.
            table_format (`Literal["markdown", "json"]`, \
             default to "markdown"):
                The format to extract table content. Note if the table cell
                contains `\n`, the Markdown format may not render correctly.
                In that case, you can use the `json` format, which extracts
                the table as a JSON string of a `list[list[str]]` object.
        """
        self._validate_init_params(chunk_size, split_by)

        if table_format not in ["markdown", "json"]:
            raise ValueError(
                "The table_format must be one of 'markdown' or 'json', "
                f"got {table_format}",
            )

        self.chunk_size = chunk_size
        self.split_by = split_by
        self.include_image = include_image
        self.separate_slide = separate_slide
        self.separate_table = separate_table
        self.table_format = table_format

        # Use TextReader to do the chunking
        self._text_reader = TextReader(self.chunk_size, self.split_by)

    def _validate_init_params(self, chunk_size: int, split_by: str) -> None:
        """Validate initialization parameters.

        Args:
            chunk_size (`int`):
                The chunk size to validate.
            split_by (`str`):
                The split mode to validate.
        """
        if chunk_size <= 0:
            raise ValueError(
                f"The chunk_size must be positive, got {chunk_size}",
            )

        if split_by not in ["char", "sentence", "paragraph"]:
            raise ValueError(
                "The split_by must be one of 'char', 'sentence' or "
                f"'paragraph', got {split_by}",
            )

    async def __call__(
        self,
        ppt_path: str,
    ) -> list[Document]:
        """Read a PowerPoint file, split it into chunks, and return a list of
        Document objects. The text, image, and table content will be returned
        in the same order as they appear in the PowerPoint presentation.

        Args:
            ppt_path (`str`):
                The input PowerPoint file path (.pptx file).

        Returns:
            `list[Document]`:
                A list of Document objects, where the metadata contains the
                chunked text, doc id and chunk id.
        """
        # Generate document ID
        doc_id = self.get_doc_id(ppt_path)

        # Load PowerPoint presentation
        try:
            from pptx import Presentation

            prs = Presentation(ppt_path)
        except ImportError as e:
            raise ImportError(
                "Please install python-pptx to use the PowerPoint reader. "
                "You can install it by `pip install python-pptx`.",
            ) from e
        except Exception as e:
            raise ValueError(
                f"Failed to read PowerPoint file {ppt_path}: {e}",
            ) from e

        # Process slides
        if self.separate_slide:
            return await self._process_slides_separately(prs, doc_id)
        else:
            return await self._process_slides_merged(prs, doc_id)

    async def _process_slides_merged(
        self,
        prs: Any,
        doc_id: str,
    ) -> list[Document]:
        """Process all slides as a merged document, maintaining order of
        text, table, and image content.

        Args:
            prs (`Any`):
                The python-pptx Presentation object.
            doc_id (`str`):
                The document ID.

        Returns:
            `list[Document]`:
                A list of Document objects from all slides merged together,
                maintaining content order.
        """
        # Get all blocks from all slides in order
        all_blocks = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks = self._get_slide_blocks(slide, slide_idx)
            all_blocks.extend(slide_blocks)

        # Convert blocks to documents
        return await self._blocks_to_documents(all_blocks, doc_id)

    async def _process_slides_separately(
        self,
        prs: Any,
        doc_id: str,
    ) -> list[Document]:
        """Process each slide as separate documents.

        Args:
            prs (`Any`):
                The python-pptx Presentation object.
            doc_id (`str`):
                The document ID.

        Returns:
            `list[Document]`:
                A list of Document objects with each slide processed
                separately.
        """
        all_docs = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks = self._get_slide_blocks(slide, slide_idx)
            slide_docs = await self._blocks_to_documents(slide_blocks, doc_id)
            all_docs.extend(slide_docs)

        return all_docs

    def _get_slide_blocks(
        self,
        slide: Any,
        slide_idx: int,
    ) -> list[TextBlock | ImageBlock]:
        """Extract all data blocks from a slide in order (text, table, image).

        Args:
            slide (`Any`):
                The slide object from python-pptx.
            slide_idx (`int`):
                The index of the slide.

        Returns:
            `list[TextBlock | ImageBlock]`:
                A list of data blocks extracted from the slide, maintaining
                the order they appear in the slide.
        """
        blocks: list[TextBlock | ImageBlock] = []
        last_type = None
        slide_header = f"[Slide {slide_idx + 1}]"

        for shape in slide.shapes:
            # Try to extract image, table, or text in order
            shape_type, extracted_data = self._extract_shape_content(
                shape,
                slide_idx,
            )

            if shape_type == "image" and extracted_data:
                # Type narrowing: when shape_type is "image",
                # extracted_data is list[ImageBlock]
                if isinstance(extracted_data, list):
                    blocks.extend(extracted_data)
                    last_type = "image"
            elif shape_type == "table" and extracted_data:
                # Type narrowing: when shape_type is "table",
                # extracted_data is str
                if isinstance(extracted_data, str):
                    last_type = self._add_table_block(
                        blocks,
                        extracted_data,
                        last_type,
                        slide_header,
                    )
            elif shape_type == "text" and extracted_data:
                # Type narrowing: when shape_type is "text", extracted_data
                # is str
                if isinstance(extracted_data, str):
                    last_type = self._add_text_block(
                        blocks,
                        extracted_data,
                        last_type,
                        slide_header,
                    )

        return blocks

    def _extract_shape_content(
        self,
        shape: Any,
        slide_idx: int,
    ) -> tuple[str | None, list[ImageBlock] | str | None]:
        """Extract content from a shape (image, table, or text).

        Args:
            shape (`Any`):
                The shape object from python-pptx.
            slide_idx (`int`):
                The index of the slide (for error logging).

        Returns:
            `tuple[str | None, list[ImageBlock] | str | None]`:
                A tuple of (content_type, content_data).
                content_type can be "image", "table", "text", or None.
        """
        # Check for images first
        if self.include_image:
            shape_images = _extract_images_from_shape(shape)
            if shape_images:
                return ("image", shape_images)

        # Check for tables
        if hasattr(shape, "has_table") and shape.has_table:
            try:
                table_data = _extract_table_data(shape.table)
                if self.table_format == "markdown":
                    return ("table", self._table_to_markdown(table_data))
                return ("table", self._table_to_json(table_data))
            except Exception as e:
                logger.warning(
                    "Failed to extract table from slide %d: %s",
                    slide_idx + 1,
                    e,
                )
                return (None, None)

        # Extract text from text frames
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            try:
                text_frame = shape.text_frame
                text_parts = [
                    para.text.strip()
                    for para in text_frame.paragraphs
                    if para.text.strip()
                ]
                if text_parts:
                    return ("text", "\n".join(text_parts))
            except Exception as e:
                logger.warning(
                    "Failed to extract text from shape in slide %d: %s",
                    slide_idx + 1,
                    e,
                )

        return (None, None)

    def _add_table_block(
        self,
        blocks: list[TextBlock | ImageBlock],
        table_text: str,
        last_type: str | None,
        slide_header: str,
    ) -> str:
        """Add a table block to the blocks list.

        Args:
            blocks (`list[TextBlock | ImageBlock]`):
                The list of blocks to add to.
            table_text (`str`):
                The formatted table text.
            last_type (`str | None`):
                The type of the last block.
            slide_header (`str`):
                The slide header to prepend if this is the first block.

        Returns:
            `str`:
                The updated last_type ("table").
        """
        should_merge = (
            not self.separate_table
            and last_type in ["text", "table"]
            and blocks
        )

        if should_merge:
            blocks[-1]["text"] += "\n" + table_text
        else:
            if last_type is None:
                table_text = slide_header + "\n" + table_text
            blocks.append(
                TextBlock(
                    type="text",
                    text=table_text,
                ),
            )

        return "table"

    def _add_text_block(
        self,
        blocks: list[TextBlock | ImageBlock],
        text: str,
        last_type: str | None,
        slide_header: str,
    ) -> str:
        """Add a text block to the blocks list.

        Args:
            blocks (`list[TextBlock | ImageBlock]`):
                The list of blocks to add to.
            text (`str`):
                The text content.
            last_type (`str | None`):
                The type of the last block.
            slide_header (`str`):
                The slide header to prepend if this is the first block.

        Returns:
            `str`:
                The updated last_type ("text").
        """
        should_merge = (
            last_type == "text"
            or (last_type == "table" and not self.separate_table)
        ) and blocks

        if should_merge:
            blocks[-1]["text"] += "\n" + text
        else:
            if last_type is None:
                text = slide_header + "\n" + text
            blocks.append(
                TextBlock(
                    type="text",
                    text=text,
                ),
            )

        return "text"

    async def _blocks_to_documents(
        self,
        blocks: list[TextBlock | ImageBlock],
        doc_id: str,
    ) -> list[Document]:
        """Convert data blocks to Document objects.

        Args:
            blocks (`list[TextBlock | ImageBlock]`):
                A list of data blocks.
            doc_id (`str`):
                The document ID.

        Returns:
            `list[Document]`:
                A list of Document objects.
        """
        documents = []

        for block in blocks:
            if block["type"] == "text":
                # Process text blocks through TextReader for chunking
                text_docs = await self._text_reader(block["text"])
                for doc in text_docs:
                    # Update doc_id but keep other metadata
                    doc.metadata.doc_id = doc_id
                    doc.id = doc_id
                    documents.append(doc)
            elif block["type"] == "image":
                # Images are independent documents
                documents.append(
                    Document(
                        metadata=DocMetadata(
                            content=block,
                            doc_id=doc_id,
                            chunk_id=0,  # Will be set later
                            total_chunks=1,
                        ),
                    ),
                )

        # Set chunk ids and total chunks
        total_chunks = len(documents)
        for idx, doc in enumerate(documents):
            doc.metadata.chunk_id = idx
            doc.metadata.total_chunks = total_chunks

        return documents

    @staticmethod
    def _table_to_markdown(table_data: list[list[str]]) -> str:
        """Convert table data to Markdown format.

        Args:
            table_data (`list[list[str]]`):
                Table data represented as a 2D list.

        Returns:
            `str`:
                Table in Markdown format.
        """
        if not table_data:
            return ""

        num_cols = len(table_data[0]) if table_data else 0
        if num_cols == 0:
            return ""

        md_table = ""

        # Header row
        header_row = "| " + " | ".join(table_data[0]) + " |\n"
        md_table += header_row

        # Separator row
        separator_row = "| " + " | ".join(["---"] * num_cols) + " |\n"
        md_table += separator_row

        # Data rows
        for row in table_data[1:]:
            # Ensure row has same number of columns as header
            while len(row) < num_cols:
                row.append("")
            data_row = "| " + " | ".join(row[:num_cols]) + " |\n"
            md_table += data_row

        return md_table

    @staticmethod
    def _table_to_json(table_data: list[list[str]]) -> str:
        """Convert table data to JSON string.

        Args:
            table_data (`list[list[str]]`):
                Table data represented as a 2D list.

        Returns:
            `str`:
                A JSON string representing the table as a 2D array,
                prefixed with a system-info tag.
        """
        json_str = json.dumps(table_data, ensure_ascii=False)
        return (
            "<system-info>A table loaded as a JSON array:</system-info>\n"
            + json_str
        )

    def get_doc_id(self, ppt_path: str) -> str:
        """Generate unique document ID from file path.

        Args:
            ppt_path (`str`):
                The path to the PowerPoint file.

        Returns:
            `str`:
                The document ID (SHA256 hash of the file path).
        """
        return hashlib.sha256(ppt_path.encode("utf-8")).hexdigest()
