# -*- coding: utf-8 -*-
"""Read tool test case."""
import base64
import os
import tempfile
from unittest.async_case import IsolatedAsyncioTestCase

from agentscope.tool import ToolChunk, Read
from agentscope.permission import (
    PermissionContext,
    PermissionBehavior,
    PermissionRule,
)
from agentscope.message import TextBlock, DataBlock, Base64Source


# pylint: disable=too-many-public-methods
class ReadToolTest(IsolatedAsyncioTestCase):
    """The read tool test case."""

    async def asyncSetUp(self) -> None:
        """The async setup method."""
        self.read_tool = Read()
        # Create a temporary file for testing
        self.temp_file = tempfile.NamedTemporaryFile(
            mode="w",
            delete=False,
            suffix=".txt",
        )
        # Write multiple lines
        for i in range(1, 11):
            self.temp_file.write(f"Line {i}\n")
        self.temp_file.close()

    async def asyncTearDown(self) -> None:
        """Clean up temporary files."""
        if os.path.exists(self.temp_file.name):
            os.unlink(self.temp_file.name)

    async def test_tool_properties(self) -> None:
        """Test read tool properties."""
        self.assertEqual(self.read_tool.name, "Read")
        self.assertIsInstance(self.read_tool.description, str)
        self.assertIsInstance(self.read_tool.input_schema, dict)
        self.assertFalse(self.read_tool.is_mcp)
        self.assertTrue(self.read_tool.is_read_only)
        self.assertTrue(self.read_tool.is_concurrency_safe)

    async def test_check_permissions(self) -> None:
        """Test read tool permission checking."""
        context = PermissionContext()
        tool_input = {"file_path": "/tmp/test.txt"}
        decision = await self.read_tool.check_permissions(tool_input, context)

        # Read/Glob/Grep are read-only, return PASSTHROUGH
        self.assertEqual(decision.behavior, PermissionBehavior.PASSTHROUGH)

    async def test_simple_read(self) -> None:
        """Test simple file reading."""
        chunk = await self.read_tool(file_path=self.temp_file.name)

        self.assertIsInstance(chunk, ToolChunk)
        self.assertEqual(chunk.state, "running")
        self.assertEqual(len(chunk.content), 1)
        self.assertIsInstance(chunk.content[0], TextBlock)

        content = chunk.content[0].text
        # Should contain all lines with line numbers
        self.assertIn("Line 1", content)
        self.assertIn("Line 10", content)

    async def test_read_with_offset(self) -> None:
        """Test reading with offset."""
        chunk = await self.read_tool(
            file_path=self.temp_file.name,
            offset=5,
        )

        self.assertEqual(chunk.state, "running")
        content = chunk.content[0].text

        # Should start from line 5
        self.assertIn("Line 5", content)
        # Line 1 should not appear (but Line 10 contains "1",
        # so check more specifically)
        lines = content.split("\n")
        line_numbers = [
            int(line.split("\t")[0].strip()) for line in lines if line.strip()
        ]
        self.assertNotIn(1, line_numbers)
        self.assertIn(5, line_numbers)

    async def test_read_with_limit(self) -> None:
        """Test reading with limit."""
        chunk = await self.read_tool(
            file_path=self.temp_file.name,
            offset=1,
            limit=3,
        )

        self.assertEqual(chunk.state, "running")
        content = chunk.content[0].text

        # Should only read 3 lines
        self.assertIn("Line 1", content)
        self.assertIn("Line 2", content)
        self.assertIn("Line 3", content)
        self.assertNotIn("Line 4", content)

    async def test_read_nonexistent_file(self) -> None:
        """Test reading a non-existent file."""
        chunk = await self.read_tool(file_path="/nonexistent/file.txt")

        self.assertEqual(chunk.state, "error")
        self.assertIn("does not exist", chunk.content[0].text)

    async def test_read_directory(self) -> None:
        """Test reading a directory (should fail)."""
        temp_dir = tempfile.mkdtemp()
        try:
            chunk = await self.read_tool(file_path=temp_dir)

            self.assertEqual(chunk.state, "error")
            self.assertIn("directory", chunk.content[0].text.lower())
        finally:
            os.rmdir(temp_dir)

    async def test_match_rule_glob_pattern(self) -> None:
        """Test match_rule with glob patterns."""
        # Test exact match
        self.assertTrue(
            await self.read_tool.match_rule(
                "test.py",
                {"file_path": "test.py"},
            ),
        )

        # Test wildcard pattern
        self.assertTrue(
            await self.read_tool.match_rule(
                "*.py",
                {"file_path": "test.py"},
            ),
        )

        # Test directory pattern
        self.assertTrue(
            await self.read_tool.match_rule(
                "/tmp/**",
                {"file_path": "/tmp/test.py"},
            ),
        )

        # Test non-matching pattern
        self.assertFalse(
            await self.read_tool.match_rule(
                "*.txt",
                {"file_path": "test.py"},
            ),
        )

        # Test empty file_path
        self.assertFalse(
            await self.read_tool.match_rule(
                "*.py",
                {"file_path": ""},
            ),
        )

    async def test_generate_suggestions(self) -> None:
        """Test generate_suggestions for file operations."""

        # Test suggestion for file in subdirectory
        suggestions = await self.read_tool.generate_suggestions(
            {"file_path": "/tmp/project/src/main.py"},
        )

        self.assertIsInstance(suggestions, list)
        self.assertGreater(len(suggestions), 0)
        self.assertIsInstance(suggestions[0], PermissionRule)

        # Should suggest parent directory pattern
        suggestion_contents = [s.rule_content for s in suggestions]
        self.assertIn("/tmp/project/src/**", suggestion_contents)

        # Test suggestion for file in root
        suggestions = await self.read_tool.generate_suggestions(
            {"file_path": "/test.py"},
        )
        self.assertGreater(len(suggestions), 0)

    async def test_read_image_file_returns_data_block(self) -> None:
        """Test reading an image file returns DataBlock."""
        img_data = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".png",
        ) as f:
            f.write(img_data)
            img_path = f.name

        try:
            chunk = await self.read_tool(file_path=img_path)

            self.assertIsInstance(chunk, ToolChunk)
            self.assertEqual(chunk.state, "running")
            self.assertEqual(len(chunk.content), 1)
            self.assertIsInstance(chunk.content[0], DataBlock)

            block = chunk.content[0]
            self.assertIsInstance(block.source, Base64Source)
            self.assertEqual(block.source.media_type, "image/png")
            self.assertEqual(block.name, os.path.basename(img_path))

            decoded = base64.b64decode(block.source.data)
            self.assertEqual(decoded, img_data)
        finally:
            os.unlink(img_path)

    async def test_read_jpeg_file_returns_data_block(self) -> None:
        """Test reading a JPEG file returns DataBlock."""
        jpg_data = b"\xff\xd8\xff\xe0" + b"\x00" * 50
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".jpg",
        ) as f:
            f.write(jpg_data)
            jpg_path = f.name

        try:
            chunk = await self.read_tool(file_path=jpg_path)

            self.assertEqual(chunk.state, "running")
            self.assertIsInstance(chunk.content[0], DataBlock)
            self.assertEqual(
                chunk.content[0].source.media_type,
                "image/jpeg",
            )
        finally:
            os.unlink(jpg_path)

    async def test_read_audio_file_returns_data_block(self) -> None:
        """Test reading an audio file returns DataBlock."""
        audio_data = b"\x00" * 200
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".mp3",
        ) as f:
            f.write(audio_data)
            mp3_path = f.name

        try:
            chunk = await self.read_tool(file_path=mp3_path)

            self.assertEqual(chunk.state, "running")
            self.assertIsInstance(chunk.content[0], DataBlock)
            self.assertEqual(
                chunk.content[0].source.media_type,
                "audio/mpeg",
            )
        finally:
            os.unlink(mp3_path)

    async def test_read_pdf_file(self) -> None:
        """Test reading a PDF file extracts text."""
        try:
            from pypdf import PdfWriter
        except ImportError:
            self.skipTest("pypdf not installed")

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pdf",
        ) as f:
            writer.write(f)
            pdf_path = f.name

        try:
            chunk = await self.read_tool(file_path=pdf_path)

            self.assertIsInstance(chunk, ToolChunk)
            self.assertEqual(chunk.state, "running")
            self.assertEqual(len(chunk.content), 1)
            self.assertIsInstance(chunk.content[0], TextBlock)
            self.assertIn("--- Page 1/1 ---", chunk.content[0].text)
        finally:
            os.unlink(pdf_path)

    async def test_read_pdf_with_pages_param(self) -> None:
        """Test reading specific pages from a PDF."""
        try:
            from pypdf import PdfWriter
        except ImportError:
            self.skipTest("pypdf not installed")

        writer = PdfWriter()
        for _ in range(5):
            writer.add_blank_page(width=612, height=792)

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pdf",
        ) as f:
            writer.write(f)
            pdf_path = f.name

        try:
            chunk = await self.read_tool(
                file_path=pdf_path,
                pages=[1, 3],
            )

            self.assertEqual(chunk.state, "running")
            text = chunk.content[0].text
            self.assertIn("--- Page 1/5 ---", text)
            self.assertIn("--- Page 3/5 ---", text)
            self.assertNotIn("--- Page 2/5 ---", text)
        finally:
            os.unlink(pdf_path)

    async def test_read_pdf_invalid_pages_filtered(self) -> None:
        """Test that out-of-range pages are filtered out."""
        try:
            from pypdf import PdfWriter
        except ImportError:
            self.skipTest("pypdf not installed")

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pdf",
        ) as f:
            writer.write(f)
            pdf_path = f.name

        try:
            chunk = await self.read_tool(
                file_path=pdf_path,
                pages=[1, 99],
            )

            self.assertEqual(chunk.state, "running")
            text = chunk.content[0].text
            self.assertIn("--- Page 1/1 ---", text)
            self.assertNotIn("Page 99", text)
        finally:
            os.unlink(pdf_path)

    async def test_read_unknown_extension_as_text(self) -> None:
        """Test that unknown extensions are read as text."""
        with tempfile.NamedTemporaryFile(
            mode="w",
            delete=False,
            suffix=".xyz",
        ) as f:
            f.write("hello world\n")
            path = f.name

        try:
            chunk = await self.read_tool(file_path=path)

            self.assertEqual(chunk.state, "running")
            self.assertIsInstance(chunk.content[0], TextBlock)
            self.assertIn("hello world", chunk.content[0].text)
        finally:
            os.unlink(path)

    async def test_read_docx_file(self) -> None:
        """Test reading a Word .docx file extracts text."""
        try:
            from docx import Document
        except ImportError:
            self.skipTest("python-docx not installed")

        doc = Document()
        doc.add_paragraph("Hello from Word")
        doc.add_paragraph("Second paragraph")

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".docx",
        ) as f:
            doc.save(f.name)
            docx_path = f.name

        try:
            chunk = await self.read_tool(file_path=docx_path)

            self.assertEqual(chunk.state, "running")
            self.assertTrue(len(chunk.content) >= 1)
            text_parts = [
                b.text for b in chunk.content if isinstance(b, TextBlock)
            ]
            combined = "\n".join(text_parts)
            self.assertIn("Hello from Word", combined)
            self.assertIn("Second paragraph", combined)
        finally:
            os.unlink(docx_path)

    async def test_read_excel_file(self) -> None:
        """Test reading an Excel file extracts table data."""
        try:
            import pandas as pd
        except ImportError:
            self.skipTest("pandas not installed")

        df = pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [25, 30]})
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".xlsx",
        ) as f:
            df.to_excel(f.name, index=False)
            xlsx_path = f.name

        try:
            chunk = await self.read_tool(file_path=xlsx_path)

            self.assertEqual(chunk.state, "running")
            text_parts = [
                b.text for b in chunk.content if isinstance(b, TextBlock)
            ]
            combined = "\n".join(text_parts)
            self.assertIn("Alice", combined)
            self.assertIn("Bob", combined)
        finally:
            os.unlink(xlsx_path)

    async def test_read_pptx_file(self) -> None:
        """Test reading a PowerPoint file extracts text."""
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Slide body text"

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pptx",
        ) as f:
            prs.save(f.name)
            pptx_path = f.name

        try:
            chunk = await self.read_tool(file_path=pptx_path)

            self.assertEqual(chunk.state, "running")
            text_parts = [
                b.text for b in chunk.content if isinstance(b, TextBlock)
            ]
            combined = "\n".join(text_parts)
            self.assertIn("Slide Title", combined)
            self.assertIn("Slide body text", combined)
        finally:
            os.unlink(pptx_path)
